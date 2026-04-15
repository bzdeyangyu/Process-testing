from __future__ import annotations

import sys
from pathlib import Path

ROOT_DIR = Path(__file__).resolve().parent
SRC_DIR = ROOT_DIR / "src"
for candidate in (ROOT_DIR, SRC_DIR):
    if str(candidate) not in sys.path:
        sys.path.insert(0, str(candidate))

from pptx import Presentation

from config import OUTPUT_DIR


def _add_title_slide(presentation: Presentation, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(presentation: Presentation, title: str, bullets: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    text_frame.text = bullets[0]
    for bullet in bullets[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0


def generate_pptx(output_path: Path | None = None) -> Path:
    presentation = Presentation()
    _add_title_slide(
        presentation,
        "AI 空间设计工作流方案",
        "基于 Skills 优先、真实 GLM API 与 MVP 三节点切片的完整交付",
    )

    slide_specs = [
        (
            "项目目标",
            [
                "面向 800 平方米科技展厅，构建可演示、可评测、可扩展的空间设计 Agent Runtime。",
                "先以 MVP 三节点跑通核心价值，再用 stub 补全 12 个 Specialist 的全链路结构。",
                "交付物覆盖 Demo、评测报告、事件日志、知识库和完整 PPT。",
            ],
        ),
        (
            "为什么做成 Runtime",
            [
                "把 Prompt、工具、规则、知识、状态机拆成稳定模块，避免业务逻辑被单次对话绑死。",
                "让后续风格库、wiki、hook、评测和真实模型接入都可以独立演进。",
                "保证单 Agent 跑通后，可以自然扩展到多 Agent 协作和前后端接入。",
            ],
        ),
        (
            "三级架构总览",
            [
                "Orchestrator 负责流程编排与共享上下文。",
                "Research / Creative / Tech / PM 四个 Leader 对应不同任务簇。",
                "12 个 Specialist 以 Skill 形式注册，当前 MVP 完整实现 3 个，其余先 stub。",
            ],
        ),
        (
            "端到端流程",
            [
                "brief 输入后先经 req_parser 形成 structured_brief。",
                "material_style 读取 DESIGN.md、生成 palette/material_spec，并通过 wiki_update 回写知识。",
                "visual_prompt 从 Agent Prompt Guide 取模板，注入项目约束，输出效果图 Prompt。",
            ],
        ),
        (
            "MVP 三节点价值",
            [
                "req_parser 验证需求结构化与工具白名单收窄。",
                "material_style 验证风格库读取和知识回写机制。",
                "visual_prompt 验证 Prompt 模板复用和真实 GLM 输出对齐。",
            ],
        ),
        (
            "Req Parser 输出",
            [
                "核心字段：project_type、audience、area_sqm。",
                "补充字段：style_preferences、special_requirements、validated。",
                "价值：为后续 Agent 提供稳定、可校验的输入契约。",
            ],
        ),
        (
            "Material Style 输出",
            [
                "核心字段：style_key、palette、materials、direction。",
                "数据来源：DESIGN.md 中的 HEX、MATERIAL 和 Prompt Guide。",
                "副作用：通过 wiki_update 把材料与风格知识沉淀到 wiki。",
            ],
        ),
        (
            "Visual Prompt 输出",
            [
                "核心字段：prompt、style_key、render_intent。",
                "Prompt 中显式注入项目类型、面积、色值和材料清单。",
                "真实 GLM 输出允许局部覆盖，fallback 保证结构稳定。",
            ],
        ),
        (
            "全量 12 阶段视图",
            [
                "Research：req_parser、case_research、concept。",
                "Creative：material_style、storyline、zoning、visual_prompt。",
                "Tech / PM：video_script、cost_estimate、report、feedback、progress。",
            ],
        ),
        (
            "Runtime 核心机制",
            [
                "ReAct 主循环：发消息、收回复、执行工具、回写结果、判终止。",
                "轮首检查 stop/cancel 信号，保证外部中断即时生效。",
                "动态重载工具列表，Post hook 新注册工具下一轮立刻可见。",
            ],
        ),
        (
            "状态机与人工干预",
            [
                "运行态覆盖 waiting_input、waiting_confirmation、running、completed、error。",
                "intervention 消息先入队，等 tool_result 完成后再统一写入，避免破坏协议。",
                "后端是唯一真相源，前端切 session 时只拉 last_messages 与 is_running。",
            ],
        ),
        (
            "工具与安全边界",
            [
                "ToolRegistry 区分原子工具与 Skill 级能力。",
                "shell/git 一律走 program + args，拒绝字符串拼接。",
                "超长输出进入外置存储并返回 blob id，防止上下文被工具输出淹没。",
            ],
        ),
        (
            "Hooks 与策略规则",
            [
                "PreToolUse 承接工具白名单、规则注入和决策拦截。",
                "PostToolUse 负责 wiki_update 等副作用沉淀。",
                "TaskCompleted / Stop hooks 用于日志收口、状态转换和取消监听。",
            ],
        ),
        (
            "事件流与观测",
            [
                "从第一行开始埋点，重点埋在 decision point 而不是执行细节。",
                "四种状态：triggered、skipped、blocked、error；没事件是更危险的信号。",
                "输出 mvp_log.json / demo_log.json，支持 JSONL sink 和内存收集。",
            ],
        ),
        (
            "评测与失败归因",
            [
                "evaluation/fixtures 提供 MVP 9 条和 full stub 用例。",
                "run_eval.py 输出 pass_rate、tool_param_accuracy、workflow_completion_rate 和 failures。",
                "消融结果至少保留一轮机制对比，避免默认新增机制一定有效。",
            ],
        ),
        (
            "Prompt 模板示例",
            [
                "Tech Showroom：强调电光青高亮、金属玻璃材质和沉浸媒体墙。",
                "Museum Narrative：强调暖色文博灯光、展柜节奏和叙事层次。",
                "Brand Experience：强调社交传播、动态媒体与戏剧化品牌体验。",
            ],
        ),
        (
            "Agent 三层分工示意",
            [
                "Orchestrator 管总流程和共享上下文。",
                "Leader 管任务簇和交接边界。",
                "Specialist/Skill 负责单点收窄执行，未来可逐个升级成完整 Agent。",
            ],
        ),
        (
            "DESIGN.md 库说明",
            [
                "共 6 种风格：tech-showroom、museum-narrative、commercial-exhibition、architectural-viz、brand-experience、cultural-space。",
                "每份 DESIGN.md 严格覆盖 9 节，其中 Agent Prompt Guide 至少 3 条模板。",
                "material_style 和 visual_prompt 直接依赖该库保证输出风格稳定。",
            ],
        ),
        (
            "Wiki 知识库架构",
            [
                "wiki/ 作为项目运行中的知识基座，区分 schema、预置页与运行期增量页。",
                "wiki_query 支持按关键词命中页面，wiki_update 负责幂等回写。",
                "demo 结束后执行 wiki_lint，检查 orphan_pages、conflicts、outdated_pages。",
            ],
        ),
        (
            "edict 可观测性说明",
            [
                "事件 envelope 统一承载 trace_id、kind、status、payload。",
                "MVP 日志目标不少于 15 条，full demo 不少于 50 条。",
                "日志既用于回放流程，也用于失败归因和后续 dashboard 对接。",
            ],
        ),
        (
            "MVP 切片说明页",
            [
                "数据流：brief -> structured_brief -> material_spec -> visual_prompt。",
                "工具流：parse_brief / validate_requirements -> get_design_spec / wiki_update -> style_inject / real LLM。",
                "验证指标：字段完整率、风格一致性、日志覆盖、知识回写成功率。",
            ],
        ),
        (
            "Skill -> Agent 演化路径",
            [
                "当前以 Skill wrapper 快速收窄边界、降低实现成本。",
                "当某节点需要更强自主决策时，再把 specialist 升级成独立 Agent。",
                "升级路径保持 schema、日志、hook、评测不变，避免重复推翻重写。",
            ],
        ),
    ]
    for title, bullets in slide_specs:
        _add_bullet_slide(presentation, title, bullets)

    target = output_path or OUTPUT_DIR / "空间设计AI工作流方案.pptx"
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))
    return target


if __name__ == "__main__":
    path = generate_pptx()
    print(path)
